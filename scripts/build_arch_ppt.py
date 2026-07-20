#!/usr/bin/env python3
"""build_arch_ppt.py — generate the IntentCam architecture deck (中文, final v3.6 architecture only).

Usage:  python3 scripts/build_arch_ppt.py
Output: docs/IntentCam-架构设计-v3.6.pptx

Content is a distilled, presentation-shaped version of ARCHITECTURE.md
(current architecture only — no history/evolution).  Re-run after
architectural changes; the deck is a generated artifact.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── palette (matches the app's accent clusters) ───────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x4F, 0x8C, 0xFF)   # accentDelegate
PINK = RGBColor(0xE6, 0x4A, 0x8C)   # accentExecute
INK = RGBColor(0x22, 0x2A, 0x38)
GRAY = RGBColor(0x5A, 0x64, 0x72)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD5, 0xDC, 0xE6)

FONT = "Noto Sans CJK SC"

SW, SH = Inches(13.333), Inches(7.5)  # 16:9

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _cjk(run):
    """Set the East-Asian typeface on a run (python-pptx sets latin only)."""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def style(run, size, color=INK, bold=False, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = FONT
    _cjk(run)


def box(slide, x, y, w, h, fill=WHITE, line=None, shadow=False, round_=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if round_:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, sz, c, b) in para:
            r = p.add_run()
            r.text = t
            style(r, sz, c, b)
    return tb


def bullets(slide, x, y, w, h, items, size=15, gap=8, color=INK, marker_color=BLUE):
    """items: list of (head, body) or plain strings."""
    paras = []
    for it in items:
        if isinstance(it, tuple):
            head, body = it
            paras.append([("▸ ", size, marker_color, True), (head, size, color, True),
                          (("  " + body) if body else "", size, GRAY, False)])
        else:
            paras.append([("▸ ", size, marker_color, True), (it, size, color, False)])
    return txt(slide, x, y, w, h, paras, space_after=gap, line_spacing=1.08)


def header(slide, title, sub=None, num=None):
    box(slide, 0, 0, SW, Inches(0.10), fill=BLUE)
    txt(slide, Inches(0.55), Inches(0.30), Inches(11.5), Inches(0.75),
        [[(title, 26, NAVY, True)]])
    if sub:
        txt(slide, Inches(0.55), Inches(0.92), Inches(12.2), Inches(0.4),
            [[(sub, 12.5, GRAY, False)]])
    if num:
        txt(slide, SW - Inches(0.9), Inches(0.30), Inches(0.55), Inches(0.4),
            [[(str(num), 12, LINE, True)]], align=PP_ALIGN.RIGHT)


def arrow(slide, x, y, w=Inches(0.42), h=Inches(0.28), color=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def down_arrow(slide, x, y, w=Inches(0.28), h=Inches(0.34), color=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def chip_node(slide, x, y, w, h, title, sub, fill=LIGHT, tcolor=NAVY, scolor=GRAY, tsize=14, ssize=10.5):
    box(slide, x, y, w, h, fill=fill, line=LINE, round_=True)
    txt(slide, x + Inches(0.12), y + Inches(0.08), w - Inches(0.24), h - Inches(0.16),
        [[(title, tsize, tcolor, True)], [(sub, ssize, scolor, False)]],
        align=PP_ALIGN.CENTER, space_after=2)


def table(slide, x, y, w, rows, col_w, header_fill=NAVY, size=12, hsize=12.5, row_h=0.42):
    gt = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w,
                                Inches(row_h * len(rows))).table
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        gt.rows[i].height = Inches(row_h)
        for j, cell in enumerate(row):
            c = gt.cell(i, j)
            c.margin_left = Inches(0.08)
            c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.03)
            c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            if i == 0:
                c.fill.fore_color.rgb = header_fill
            else:
                c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            p = c.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(cell)
            style(r, hsize if i == 0 else size,
                  WHITE if i == 0 else INK, i == 0)
    return gt


# ═════════════════════════════ S1 · 封面 ═══════════════════════════
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, SW, SH, fill=NAVY)
box(s, 0, SH - Inches(0.18), SW, Inches(0.18), fill=BLUE)
box(s, Inches(0.9), Inches(2.05), Inches(0.16), Inches(1.7), fill=BLUE)
txt(s, Inches(1.25), Inches(2.0), Inches(10.8), Inches(1.9),
    [[("IntentCam 意图相机", 44, WHITE, True)],
     [("关键架构设计", 26, RGBColor(0xBFD-0x100 if False else 0x9D, 0xC4, 0xFF), False)]], space_after=10)
txt(s, Inches(1.25), Inches(4.15), Inches(10.5), Inches(1.6),
    [[("拍一张照片 → 多轮 LLM 识别 → 可执行的下一步(action chips)", 16, RGBColor(0xD8, 0xE2, 0xF2), False)],
     [("v3.6 最终架构 · 2026-07", 13, RGBColor(0x8E, 0x9C, 0xB8), False)]], space_after=8)

# ═════════════════════ S2 · 产品定位 ═══════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "产品定位:照片 = 意图的入口", "What can the user do with this frame?", 2)
bullets(s, Inches(0.55), Inches(1.5), Inches(6.6), Inches(4.6), [
    ("一句话", "对取景内容做一次快门,得到结构化 bubble + 可点动作 chips,而不是一段文字描述。"),
    ("动作优先 (action-first)", "系统围绕「动作」组织:LLM 直接提议 action_ids,无意图分类法。"),
    ("端侧先行", "双 JPEG + 端侧 OCR hint 先行,LLM 只做理解、决策与转录。"),
    ("可验证", "每个动作都有 required input;缺参不终结,nudge 再查一轮。"),
], size=14.5, gap=14)
box(s, Inches(7.5), Inches(1.5), Inches(5.3), Inches(4.9), fill=LIGHT, line=LINE, round_=True)
txt(s, Inches(7.8), Inches(1.7), Inches(4.8), Inches(0.5), [[("一次快门发生的事", 15, NAVY, True)]])
steps = [("① 帧采集", "3200px 缩略 + 4096px 原图 + 端侧 OCR"),
         ("② 多轮识别", "≤4 轮工具调用,zoom_in 细看局部"),
         ("③ 结构化输出", "bubble:摘要 + 细节表 + action_ids"),
         ("④ 动作执行", "地图/拨号/分享/标签页 一触即达")]
for i, (t, d) in enumerate(steps):
    y = Inches(2.25) + i * Inches(1.02)
    chip_node(s, Inches(7.8), y, Inches(4.7), Inches(0.82), t, d, fill=WHITE)

# ═════════════════ S3 · 端到端架构总览 ════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "端到端架构总览", "快门 → bubble → chips,一条有界流水线", 3)
chain = [("FrameAnalyzer", "双 JPEG\n+ OCR hint"), ("CycleManager", "队列 8\nworker 2"),
         ("ToolUseLoop", "≤4 轮\n4 工具"), ("Bubble", "摘要+细节\n+action_ids"),
         ("动作解析", "提议∩注册∩启用\n+rescue"), ("执行", "Intent/Toast\n/表单/页面")]
x = Inches(0.45)
for i, (t, d) in enumerate(chain):
    chip_node(s, x, Inches(1.85), Inches(1.82), Inches(1.25), t, d,
              fill=LIGHT if i not in (3, 5) else RGBColor(0xE8, 0xF0, 0xFF), tsize=13.5)
    if i < len(chain) - 1:
        arrow(s, x + Inches(1.86), Inches(2.32))
    x += Inches(2.17)
txt(s, Inches(0.45), Inches(3.45), Inches(12.4), Inches(0.4),
    [[("两层边界:UI 只读 StateFlow;LLM 只通过工具与契约字段交互 —— 两侧都可独立替换。", 13, GRAY, False)]])
table(s, Inches(0.45), Inches(3.95), Inches(12.4), [
    ["层", "职责", "关键约束"],
    ["采集层", "CameraX 帧 → 双 JPEG + 端侧 OCR hint", "缩略图给 LLM,原图只服务 zoom_in 裁剪"],
    ["编排层", "CycleManager 排队/并发/状态流", "队列 8 · 并发 2 · 90s 软超时,背压拒绝快门"],
    ["识别层", "ToolUseLoop 多轮工具调用", "temperature 0 · MAX_TOKENS 8192 · ≤4 轮"],
    ["动作层", "注册表 + 解析器 + 校验门 + rescue", "LLM 提议是唯一路由信号,rescue 只加不减"],
    ["呈现层", "Compose bubble 卡片 + chips + 整页页面", "chip 三态:Validated / Spinner / Ghost"],
], [1.6, 5.6, 5.2], row_h=0.5)

# ═════════════════ S4 · 双模块划分 ════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "双模块:Android 壳 + 纯 JVM 核", "识别管线在 :shared —— 产品代码与评测代码是同一份", 4)
box(s, Inches(0.55), Inches(1.55), Inches(5.95), Inches(4.15), fill=LIGHT, line=LINE, round_=True)
txt(s, Inches(0.85), Inches(1.75), Inches(5.4), Inches(0.5), [[(":app(Android)", 17, BLUE, True)]])
bullets(s, Inches(0.85), Inches(2.3), Inches(5.4), Inches(3.2), [
    "Compose UI(相机/卡片/详情页/标签页/设置)",
    "CameraX 帧采集 · HMS ML Kit 端侧 OCR",
    "ActionRegistry + action body(Intent/页面)",
    "设置与持久化(SharedPreferences)",
], size=13.5, gap=10)
box(s, Inches(6.85), Inches(1.55), Inches(5.95), Inches(4.15), fill=RGBColor(0xE8, 0xF0, 0xFF), line=LINE, round_=True)
txt(s, Inches(7.15), Inches(1.75), Inches(5.4), Inches(0.5), [[(":shared(纯 JVM Kotlin)", 17, NAVY, True)]])
bullets(s, Inches(7.15), Inches(2.3), Inches(5.4), Inches(3.2), [
    "ToolUseLoop / LlmClient(SSE)/ 数据模型",
    "InputParsers · ActionRescue · LabelHtml",
    "不含任何 android.* 引用",
    ("同一份类跑两端", "app 运行 + JVM eval 评测,零漂移"),
], size=13.5, gap=10)
box(s, Inches(0.55), Inches(6.0), Inches(12.25), Inches(1.0), fill=NAVY, round_=True)
txt(s, Inches(0.85), Inches(6.14), Inches(11.7), Inches(0.75),
    [[("为什么:", 13, BLUE, True),
      ("eval 不复刻识别逻辑 —— 测的就是线上代码;Android 接缝(OCR/图像/注册表)以注入方式替换。", 13, WHITE, False)]])

# ═════════════════ S5 · 帧管线与 OCR hint ═════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "帧管线:双 JPEG + OCR hint", "端侧先做能做的事,LLM 专注理解", 5)
table(s, Inches(0.55), Inches(1.55), Inches(12.25), [
    ["产物", "规格", "用途"],
    ["缩略图 JPEG", "3200px · q90", "发给 LLM(视觉上限内)+ bubble 展示"],
    ["原图 JPEG", "4096px · q95", "zoom_in 裁剪源,保住小字细节"],
    ["OCR hint", "HMS ML Kit 离线(中/英)", "第 1 轮随图注入:文字行 + 归一化 bbox"],
], [2.4, 3.2, 6.65], row_h=0.52)
bullets(s, Inches(0.55), Inches(3.85), Inches(12.2), Inches(2.6), [
    ("bbox 贯通", "OCR hint 的 bbox 供 zoom_in 定位,details[].bbox 回填详情页高亮点。"),
    ("无 OCR 工具", "OCR 自动跑:第 1 轮全图 + 每次 zoom_in 裁剪后,模型永远读得到刚放大区域的文字。"),
    ("端侧优先", "图像缩放/裁剪/编码全在端侧,网络只承担 LLM 推理。"),
], size=14, gap=12)

# ═════════════════ S6 · CycleManager ══════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "CycleManager:有界生产者/消费者", "快门是生产者,worker 池是消费者,UI 只订阅状态流", 6)
for i, (t, d, c) in enumerate([
        ("CYCLE_QUEUE_DEPTH = 8", "排队+在途上限,满了拒拍(背压)", BLUE),
        ("CYCLE_CONCURRENCY = 2", "并发 LLM 流上限(API 负载/内存)", PINK),
        ("CYCLES_MAX_TOTAL = 8", "终态任务 FIFO 淘汰", NAVY)]):
    box(s, Inches(0.55) + i * Inches(4.2), Inches(1.6), Inches(3.95), Inches(1.15), fill=LIGHT, line=c, round_=True)
    txt(s, Inches(0.75) + i * Inches(4.2), Inches(1.74), Inches(3.6), Inches(0.9),
        [[(t, 14, c, True)], [(d, 11.5, GRAY, False)]], space_after=2)
bullets(s, Inches(0.55), Inches(3.1), Inches(12.2), Inches(3.4), [
    ("状态机", "PENDING → IN_FLIGHT → COMPLETE / SUPERSEDED / ERRORED;新快门把旧周期降级为 SUPERSEDED(继续跑完不浪费)。"),
    ("CycleSnapshot", "每个周期的 status/bubble/nRounds/pendingInputs 都是独立 StateFlow —— 单周期更新不触发整表重组。"),
    ("busy 派生流", "busy = 聚焦周期处于 PENDING/IN_FLIGHT(flatMapLatest 派生,无手工标志位)。"),
    ("90s 软超时", "llmTimeoutMs 兜底挂起的 API 调用:周期标 ERRORED,保留最后可用 bubble。"),
], size=14, gap=12)

# ═════════════════ S7 · 多轮工具调用协议 ══════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "多轮工具调用:LLM 的探索-收敛回路", "temperature 0 · SSE 流式 · ≤4 轮 · emit_bubble 终止", 7)
table(s, Inches(0.55), Inches(1.5), Inches(12.25), [
    ["工具", "职责", "何时调用"],
    ["zoom_in", "按 bbox 从 4096 原图裁剪并回传(自动 OCR)", "需要看清局部/小字"],
    ["compare_text", "模型读数 vs OCR hint 的端侧 diff", "与 OCR 行不一致时校验"],
    ["extract_text", "对某区域单跑高保真 OCR,只回文字", "只要准确字符,不需再看图"],
    ["emit_bubble", "终止轮:产出结构化 bubble(必须调用)", "理解完成"],
], [2.2, 6.0, 4.05], row_h=0.5)
bullets(s, Inches(0.55), Inches(4.15), Inches(12.2), Inches(2.3), [
    ("emit_bubble 契约", "content + intent(≤30字自由文本)+ details[](kind/label/value/bbox)+ confidence + action_ids + label_markdown?"),
    ("action_ids 无枚举", "可用动作拼进系统提示 actions ∈ {…},schema 不设 enum —— 注册表单一来源,永不漂移。"),
    ("确定性", "temperature=0 + MAX_TOKENS 8192,评测可复现。"),
], size=14, gap=12)

# ═════════════════ S8 · Bubble 数据模型 ═══════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Bubble:一次识别的结构化答案", "跨模块(shared)纯数据类,UI 与评测共享同一模型", 8)
rows = [
    ("title / detail", "意图短语(≤30字)/ 场景详述,来自 emit_bubble.intent / content"),
    ("details[]", "抽取明细 kind/label/value/bbox → 详情页表格 + 位置高亮"),
    ("confidence", "0.0~1.0 置信度"),
    ("llmProposedActions", "LLM 的 action_ids 原样 —— 唯一路由信号"),
    ("actions", "解析后可见 chip 集(提议 ∩ 注册 ∩ 启用 ∪ rescue)"),
    ("validatedInputs / pendingInputs", "每个动作的必填项校验结果 → chip 三态与 nudge"),
    ("labelMarkdown?", "标签全文转录(view_label 的载荷与 rescue 线索)"),
]
table(s, Inches(0.55), Inches(1.5), Inches(12.25),
      [["字段", "含义"]] + [[a, b] for a, b in rows], [3.6, 8.65], row_h=0.5)
txt(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(0.8),
    [[("UiState 另持:phase / cycles / error / debugEnabled / pendingAction / pendingConfirmation / renderedLabel。", 12.5, GRAY, False)]])

# ═════════════════ S9 · 动作三级解析 ══════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "动作系统:三级解析 + 校验门 + rescue", "LLM 提议是唯一路由信号,系统负责可执行性", 9)
funnel = [("LLM 提议 action_ids", "模型按 prompt 映射规则直选", RGBColor(0xE8, 0xF0, 0xFF)),
          ("∩ 注册表白名单", "未知 id 丢弃(防幻觉)", LIGHT),
          ("∩ 用户启用集", "开发阶段全量启用", LIGHT),
          ("∪ content rescue", "电话/证件/收款码/标签 —— 只加不减", RGBColor(0xFB, 0xE8, 0xF1)),
          ("Bubble.actions(可见 chips)", "", RGBColor(0xE6, 0xF7, 0xEC))]
for i, (t, d, f) in enumerate(funnel):
    y = Inches(1.55) + i * Inches(0.98)
    chip_node(s, Inches(0.55), y, Inches(5.6), Inches(0.8), t, d, fill=f, tsize=13.5)
    if i < len(funnel) - 1:
        down_arrow(s, Inches(3.2), y + Inches(0.8))
box(s, Inches(6.6), Inches(1.55), Inches(6.2), Inches(4.75), fill=LIGHT, line=LINE, round_=True)
txt(s, Inches(6.9), Inches(1.75), Inches(5.6), Inches(0.5), [[("ActionOrchestrator 校验门", 15, NAVY, True)]])
bullets(s, Inches(6.9), Inches(2.3), Inches(5.65), Inches(3.8), [
    ("requiredInputs", "每个动作声明必填输入(如 dial→phone_number,view_label→label_markdown)。"),
    ("validateInputs", "对每条 emit_bubble 跑 parser,写 validatedInputs/pendingInputs。"),
    ("shouldFinalize", "缺参 → CONTINUE:注入「还缺:X、Y」nudge 再来一轮(有界重试);齐 → FINALIZE。"),
    ("rescue 时机", "校验前注入,门看到的与 UI 渲染的是同一个 enriched bubble。"),
], size=12.5, gap=8)

# ═════════════════ S10 · 六个动作 ════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "六个动作:注册表驱动", "加一个动作 = 注册一个 ActionDef + prompt 一句映射 + eval 镜像", 10)
table(s, Inches(0.55), Inches(1.5), Inches(12.25), [
    ["id", "chip", "出口", "必填输入", "簇"],
    ["open_in_maps", "在地图中打开", "geo: Intent + 选择器(可按名称搜)", "query", "DELEGATE"],
    ["dial_number", "拨号", "ACTION_DIAL(拨号器内再确认)", "phone_number", "EXECUTE"],
    ["share", "分享文本", "ACTION_SEND text/plain", "text", "DELEGATE"],
    ["view_label", "查看标签", "整页渲染页(分享图片/文字)", "label_markdown", "DELEGATE"],
    ["scan_to_pay", "扫码支付", "安全提示 Toast(永不自动起支付)", "—", "EXECUTE"],
    ["redact_id", "遮挡证件号", "提示 Toast", "—", "EXECUTE"],
], [2.3, 1.9, 4.6, 2.15, 1.3], row_h=0.52)
txt(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(1.2),
    [[("开发阶段(2026-07):", 13, PINK, True),
      ("确认弹窗与默认关闭开关已解除,chip 常显直达;机制(requiresConfirmation/userPrefKey)休眠保留,面向最终用户构建前按动作恢复。", 13, INK, False)]])

# ═════════════════ S11 · Chip 状态机 + Outcome 分发 ══════════════
s = prs.slides.add_slide(BLANK)
header(s, "Chip 三态 与 ActionOutcome 五路分发", "编译期强制的两类封闭集合", 11)
for i, (t, d, c) in enumerate([
        ("Validated", "输入齐备 · 实色可点", RGBColor(0xE6, 0xF7, 0xEC)),
        ("Spinner", "周期未终结且未校验(含已算出 false):nudge 轮可能补齐,不可点", RGBColor(0xFF, 0xF3, 0xD6)),
        ("Ghost", "周期终结仍缺输入 · 灰色可点,body 自解释(「未发现可拨打的号码」)", LIGHT)]):
    chip_node(s, Inches(0.55), Inches(1.55) + i * Inches(1.05), Inches(5.9), Inches(0.88), t, d, fill=c, tsize=13.5)
box(s, Inches(6.85), Inches(1.55), Inches(5.95), Inches(4.4), fill=LIGHT, line=LINE, round_=True)
txt(s, Inches(7.15), Inches(1.75), Inches(5.4), Inches(0.5), [[("ActionOutcome(sealed,穷举 when)", 14.5, NAVY, True)]])
bullets(s, Inches(7.15), Inches(2.35), Inches(5.45), Inches(3.4), [
    ("None", "纯 UI 副作用(预留)"),
    ("LaunchAndroidIntent", "地图/拨号/分享"),
    ("ShowUiFeedback", "Toast"),
    ("RequestArgs", "缺参弹表单,回填后再调 body"),
    ("ShowRenderedLabel", "携带副本的标签页载荷 → 全屏页面 overlay"),
], size=13, gap=8)
txt(s, Inches(0.55), Inches(5.05), Inches(5.9), Inches(1.4),
    [[("要点:", 13, BLUE, True), ("chip 可点性由「周期状态 × 校验结果」决定;新增 outcome 会强制编译器指出所有需处理的分发点。", 13, INK, False)]])

# ═════════════════ S12 · view_label 专题 ═════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "view_label:标签识别 → 整页渲染 → 分享", "零新增依赖 · 零存储权限 · 一个 FileProvider", 12)
flow = [("label_markdown 契约", "LLM 用 markdown 逐字转录整个标签\n(标题/字段列表/GFM 表格)"),
        ("LabelHtml", "手写 markdown 子集→HTML\nescape-first 防注入 · 内置 CSS 模板"),
        ("LabelPageScreen", "WebView 全屏渲染\nviewport=device-width → CSS px=dp"),
        ("整页截图", "enableSlowWholeDocumentDraw\nresize→draw→同帧还原+空白守卫"),
        ("分享", "图片:FileProvider+ClipData\n文字:ACTION_SEND text/plain")]
y = Inches(1.55)
for i, (t, d) in enumerate(flow):
    chip_node(s, Inches(0.55), y, Inches(4.35), Inches(0.95), t, d, fill=LIGHT if i % 2 else RGBColor(0xE8, 0xF0, 0xFF), tsize=12.5, ssize=10)
    if i < len(flow) - 1:
        down_arrow(s, Inches(2.55), y + Inches(0.95), h=Inches(0.26))
    y += Inches(1.22)
box(s, Inches(5.35), Inches(1.55), Inches(7.45), Inches(5.35), fill=LIGHT, line=LINE, round_=True)
txt(s, Inches(5.65), Inches(1.75), Inches(6.9), Inches(0.5), [[("关键工程决策", 15, NAVY, True)]])
bullets(s, Inches(5.65), Inches(2.3), Inches(6.9), Inches(4.4), [
    ("契约即载荷", "label_markdown 同时是渲染源与「分享文字」内容;也是 view_label 的必填输入与 rescue 线索。"),
    ("不引三方库", "markdown→HTML 转换器 ~200 行纯 Kotlin(shared,可 JVM 冒烟测试)。"),
    ("截屏上可见的 WebView", "离屏 WebView 不保证光栅化(两次空白教训);enableSlowWholeDocumentDraw 让折叠线以下内容可被绘制。"),
    ("payload 带副本", "bubble 被逐出历史后页面照常显示。"),
    ("DEBUG 钩子", "am start --ez dev_label_page true:无相机/无 LLM 全链路自检。"),
], size=12.5, gap=9)

# ═════════════════ S13 · 设置与持久化 ════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "设置与持久化:默认即可用,修改才接管", "SharedPreferences + BuildConfig 烘焙缺省 token", 13)
bullets(s, Inches(0.55), Inches(1.6), Inches(12.2), Inches(4.8), [
    ("烘焙缺省", "secrets.properties → BuildConfig.DEFAULT_AUTH_TOKEN;prefs 无 token 键时回落到它(随新 APK 轮换)。"),
    ("脏检查保存", "离开设置页仅当用户改过 LLM 字段才落盘 —— 无修改访问零写入,缺省 token 长期有效;一旦修改,用户接管整个配置。"),
    ("调试日志开关", "位于设置页,默认关闭;相机顶栏因此只剩设置齿轮。"),
    ("界面即设置", "设置页 = 接口三字段(baseUrl/token/model)+ 调试开关 + 关于;返回即保存(带脏检查),无保存/恢复默认按钮。"),
], size=14.5, gap=14)

# ═════════════════ S14 · Eval 架构 ═══════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Eval:与产品同源的度量体系", "JVM 直跑 :shared,无需设备", 14)
box(s, Inches(0.55), Inches(1.55), Inches(5.8), Inches(2.1), fill=NAVY, round_=True)
txt(s, Inches(0.85), Inches(1.8), Inches(5.2), Inches(1.6),
    [[("composite = 0.55·r_actions + 0.30·r_text + 0.15·r_inputs", 15.5, WHITE, True)],
     [("r_actions = 期望动作集合召回 · r_text = 关键词混合匹配 · r_inputs = 必填项满足率", 12, RGBColor(0xC9, 0xD6, 0xEA), False)]],
    space_after=8)
bullets(s, Inches(0.55), Inches(4.0), Inches(5.8), Inches(2.6), [
    ("镜像而非复刻", "EvalRunner 复用 ActionRescue/InputParsers,动作词表与注册表 lockstep。"),
    ("套件即回归网", "baselines.json 登记基线,|Δ|≥0.05 触发告警。"),
    ("精度哨兵", "none 套件 over_fire_rate 监控误开火。"),
], size=13, gap=9)
table(s, Inches(6.75), Inches(1.55), Inches(6.05), [
    ["套件", "基线", "n"],
    ["dial_number", "0.9049", "29"],
    ["open_in_maps", "0.9144", "34"],
    ["share", "0.9587", "38"],
    ["scan_to_pay", "0.8444", "6"],
    ["none(过火检查)", "0.9474", "16"],
    ["view_label", "0.6711", "30"],
], [2.6, 1.9, 1.55], row_h=0.48)
txt(s, Inches(6.75), Inches(5.1), Inches(6.0), Inches(0.5),
    [[("kimi k3 @ 2026-07-19 基线", 11.5, GRAY, False)]])

# ═════════════════ S15 · 关键设计决策 ════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "关键设计决策", "六条贯穿全局的取舍", 15)
decisions = [
    ("LLM 权威,无分类法", "intent 是自由文本,action_ids 是唯一路由信号 —— 没有第二种真相需要同步。"),
    ("字符串注册表,而非枚举", "action = 注册一条 ActionDef;prompt 词表运行时拼接,schema/文档/评测同一来源。"),
    ("识别逻辑只写一遍", ":shared 纯 JVM,app 与 eval 共用 —— 从根上消除「评测复刻漂移」。"),
    ("rescue 只加不减", "可验证内容线索(电话/证件/收款码/标签字段)兜底补 chip,永远不替 LLM 做减法。"),
    ("可执行性前置", "requiredInputs + 校验门:缺参不终结、chip 不可点,失败发生在用户点击之前。"),
    ("零依赖渲染/导出", "手写 markdown→HTML + 可见 WebView 截图:无三方库、无存储权限、截图路径经真机+模拟器验证。"),
]
for i, (t, d) in enumerate(decisions):
    col, row = i % 2, i // 2
    x = Inches(0.55) + col * Inches(6.3)
    y = Inches(1.6) + row * Inches(1.72)
    box(s, x, y, Inches(6.0), Inches(1.5), fill=LIGHT, line=LINE, round_=True)
    txt(s, x + Inches(0.22), y + Inches(0.14), Inches(5.6), Inches(1.25),
        [[(f"{i+1}. {t}", 14.5, NAVY, True)], [(d, 12, GRAY, False)]], space_after=4)
box(s, 0, SH - Inches(0.55), SW, Inches(0.55), fill=NAVY)
txt(s, Inches(0.55), SH - Inches(0.5), Inches(12.2), Inches(0.4),
    [[("IntentCam v3.6 · 架构即代码:registry / contract / pipeline 三处入手即可扩展。", 12, RGBColor(0xC9, 0xD6, 0xEA), False)]])

# ── save ──────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(__file__), "..", "docs", "IntentCam-架构设计-v3.6.pptx")
out = os.path.abspath(out)
prs.save(out)
print(f"saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
